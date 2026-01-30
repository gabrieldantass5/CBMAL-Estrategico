#!/usr/bin/env python3
import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Configurações de Marca (CBMAL)
COLOR_VERMELHO = RGBColor(193, 10, 10)  # Vermelho Assinatura #C10A0A
COLOR_AZUL = RGBColor(74, 148, 255)     # Azul #4A94FF
FONT_TITULO = 'Exo 2 Black'
FONT_CORPO = 'Exo 2'

# Caminhos Oficiais
PATH_BRASAO = r'c:\Users\D_A_N\OneDrive\Desktop\Planejamento Estratégico - CBMAL\99_Base_de_Conhecimento\05_Templates_Institucionais\Identidade Visual\Brasão Manual Idenditade 2022 - Sem fundo.png'

class PPTGenerator:
    def __init__(self, title, subtitle="SEÇÃO DE GESTÃO ESTRATÉGICA - APO/EMG"):
        self.prs = Presentation()
        # Ajustar para 16:9
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.add_title_slide(title, subtitle)

    def _add_decorations(self, slide, slide_number=None):
        """Adiciona o Brasão e a numeração em cada slide."""
        # Adicionar Brasão (Top Right)
        if os.path.exists(PATH_BRASAO):
            left = self.prs.slide_width - Inches(1.2)
            top = Inches(0.2)
            slide.shapes.add_picture(PATH_BRASAO, left, top, height=Inches(1.0))

        # Adicionar Número do Slide (Bottom Right)
        if slide_number:
            txBox = slide.shapes.add_textbox(self.prs.slide_width - Inches(0.8), 
                                           self.prs.slide_height - Inches(0.5), 
                                           Inches(0.5), Inches(0.3))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = str(slide_number)
            p.font.name = FONT_CORPO
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(128, 128, 128)

    def add_title_slide(self, title, subtitle):
        slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]

        title_shape.text = title.upper()
        subtitle_shape.text = subtitle

        # Estilização
        for run in title_shape.text_frame.paragraphs[0].runs:
            run.font.name = FONT_TITULO
            run.font.size = Pt(54)
            run.font.color.rgb = COLOR_VERMELHO

        self._add_decorations(slide, 1)

    def add_content_slide(self, title, content_list, slide_index):
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        body_shape = slide.placeholders[1]

        title_shape.text = title.upper()
        tf = body_shape.text_frame
        tf.word_wrap = True

        if content_list:
            # Primeiro item
            p = tf.paragraphs[0]
            p.text = content_list[0]
            p.font.name = FONT_CORPO
            p.font.size = Pt(24)

            # Demais itens
            for i in range(1, len(content_list)):
                p = tf.add_paragraph()
                p.text = content_list[i]
                p.level = 0
                p.font.name = FONT_CORPO
                p.font.size = Pt(24)

        # Estilização do Título
        for run in title_shape.text_frame.paragraphs[0].runs:
            run.font.name = FONT_TITULO
            run.font.size = Pt(36)
            run.font.color.rgb = COLOR_AZUL

        self._add_decorations(slide, slide_index)

    def save(self, filename):
        self.prs.save(filename)
        print(f"✅ Apresentação salva: {filename}")

def parse_markdown_to_ppt(md_path):
    if not os.path.exists(md_path):
        print("❌ Arquivo não encontrado.")
        return

    with open(md_path, 'r', encoding='utf-8') as f:
        lines = f.readlines()

    ppt = None
    current_title = ""
    current_content = []
    slide_count = 1

    for line in lines:
        line = line.strip()
        if line.startswith("# "):
            if ppt:
                slide_count += 1
                ppt.add_content_slide(current_title, current_content, slide_count)
            ppt = PPTGenerator(line[2:])
            current_title = ""
            current_content = []
            slide_count = 1
        elif line.startswith("## "):
            if current_title:
                slide_count += 1
                ppt.add_content_slide(current_title, current_content, slide_count)
            current_title = line[3:]
            current_content = []
        elif line.startswith("- ") or line.startswith("* "):
            current_content.append(line[2:])
        elif line and not line.startswith("!"):
            current_content.append(line)

    if ppt and current_title:
        slide_count += 1
        ppt.add_content_slide(current_title, current_content, slide_count)

    output_path = md_path.replace(".md", ".pptx")
    ppt.save(output_path)

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Uso: python scripts/ppt_generator.py <arquivo.md>")
    else:
        parse_markdown_to_ppt(sys.argv[1])
